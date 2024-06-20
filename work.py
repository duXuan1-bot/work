import jieba
import requests
import streamlit as st
from streamlit_echarts import st_echarts
from collections import Counter
import re
import string
import pandas as pd
from bs4 import BeautifulSoup
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from io import StringIO
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import os

# 读取并处理文本文件
def read_text_file(file):
    return file.read().decode("utf-8")

# 读取docx文件
def read_docx(file):
    doc = Document(file)
    full_text = [para.text for para in doc.paragraphs]
    return '\n'.join(full_text)

# 读取pptx文件
def read_pptx(file):
    prs = Presentation(file)
    full_text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
    return '\n'.join(full_text)

# 读取pdf文件
def read_pdf(file):
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")
    full_text = [page.get_text() for page in pdf_document]
    return '\n'.join(full_text)

# 数据清洗函数
def clean_text(text):
    return text.replace('\n', '').replace(' ', '').strip()

# 分词函数
def segment(text):
    stopwords = set(['的', '了', '在', '是', '我', '你', '他', '她', '它', '们', '这', '那', '之', '与', '和', '或', '等', '虽然', '大', '更', '多', '但是', '然而', '因此', '条', '年'])
    punctuation = "、，。！？；：　“”‘’~@#￥%……&*（）【】｛｝+-*/=《》<>「」『』【】〔〕｟｠«»“”‘’'':;,/\\|[]{}()$^↑"
    text = text.translate(str.maketrans("", "", punctuation)).replace('\n', '')
    words = [word for word in jieba.lcut(text) if word not in stopwords and len(word) > 1]
    return words

# 去除标点符号
def remove_punctuation(text):
    punctuation = string.punctuation
    return re.sub(r'[{}]'.format(punctuation), '', text)

# 从HTML中提取正文文本
def extract_body_text(html):
    soup = BeautifulSoup(html, 'html.parser')
    text = soup.find('body').get_text()
    return text

# 生成词云图
def generate_wordcloud(text):
    font_path = os.path.join(os.path.dirname(__file__), 'SIMHEI.TTF')
    wordcloud = WordCloud(font_path=font_path, width=800, height=400, background_color='white').generate(text)
    return wordcloud

# 统计词频并生成图表
def plot_word_frequency(word_counts):
    top_words = word_counts.most_common(20)

    # 条形图
    bar_options = {
        "tooltip": {"trigger": 'item', "formatter": '{b} : {c}'},
        "xAxis": [{"type": "category", "data": [word for word, count in top_words], "axisLabel": {"interval": 0, "rotate": 30}}],
        "yAxis": [{"type": "value"}],
        "series": [{"type": "bar", "data": [count for word, count in top_words]}]
    }
    st_echarts(bar_options, height='500px')

    # 饼图
    pie_options = {
        "tooltip": {"trigger": 'item', "formatter": '{b} : {c} ({d}%)'},
        "series": [{
            "name": '词频',
            "type": 'pie',
            "radius": '55%',
            "data": [{"value": count, "name": word} for word, count in top_words],
            "emphasis": {
                "itemStyle": {
                    "shadowBlur": 10,
                    "shadowOffsetX": 0,
                    "shadowColor": 'rgba(0, 0, 0, 0.5)'
                }
            }
        }]
    }
    st_echarts(pie_options, height='500px')

    # 添加词频排名表格
    df = pd.DataFrame(top_words, columns=['Word', 'Frequency'])
    st.write(df)

# 主函数
def run():
    st.set_page_config(page_title="欢迎使用中文文本分析工具！", page_icon="")

    st.markdown("<h1 style='font-size: 40px;'>欢迎使用中文文本分析工具！</h1>", unsafe_allow_html=True)

    # 用户选择数据输入方式
    option = st.selectbox('选择数据输入方式', ('输入 URL', '上传文本文件', '上传文档文件'))

    text = ""

    if option == '输入 URL':
        url = st.text_input('输入 URL:')
        if url:
            try:
                r = requests.get(url)
                if r.status_code == 200:
                    r.encoding = 'utf-8'
                    text = extract_body_text(r.text)
                else:
                    st.error(f"无法获取URL内容: 状态码 {r.status_code}")
            except requests.RequestException as e:
                st.error(f"无法获取URL内容: {e}")

    elif option == '上传文本文件':
        uploaded_file = st.file_uploader("上传文本文件", type=["txt"])
        if uploaded_file is not None:
            text = read_text_file(uploaded_file)

    else:
        uploaded_file = st.file_uploader("上传文档文件", type=["docx", "pptx", "pdf"])
        if uploaded_file is not None:
            if uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text = read_docx(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = read_pptx(uploaded_file)
            elif uploaded_file.type == "application/pdf":
                text = read_pdf(uploaded_file)

    if text:
        # 文本清洗和分词
        text = remove_punctuation(text)
        text = clean_text(text)
        words = segment(text)
        word_counts = Counter(words)

        # 展示最常见的词
        most_common_word, most_common_count = word_counts.most_common(1)[0]
        st.write(f"最常见的词是 “{most_common_word}” ，出现次数为 {most_common_count}次。")

        # 统计词频并生成图表
        plot_word_frequency(word_counts)

        # 生成词云
        wordcloud = generate_wordcloud(' '.join(words))
        plt.imshow(wordcloud, interpolation='bilinear')
        plt.axis("off")
        st.pyplot(plt)

if __name__ == "__main__":
    run()
